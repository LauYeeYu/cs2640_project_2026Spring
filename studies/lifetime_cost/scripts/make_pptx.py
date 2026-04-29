"""Build the AdaptiveCache decks (.pptx).

Outputs two files:
  reports/AdaptiveCache_talk.pptx      — main talk: method walkthrough, ~5 min
  reports/AdaptiveCache_appendix.pptx  — Q&A backup: project arc, vision, results

Main talk (16 slides, ~5 min):
  1.  Title
  2.  Setup — what compaction is and why we have to do it
  3.  Family tree — the 6 design families on one map
  4–15. The 12 methods, each with result + N runs callout
  16. Takeaway

Palette: navy + cream + warm orange + green (matches the figure palette).
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "reports" / "figures"
OUT_TALK = ROOT / "reports" / "AdaptiveCache_talk.pptx"
OUT_APPENDIX = ROOT / "reports" / "AdaptiveCache_appendix.pptx"

# 16:9 widescreen
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Palette
NAVY        = RGBColor(0x1E, 0x3A, 0x8A)
NAVY_DEEP   = RGBColor(0x14, 0x25, 0x5C)
CREAM       = RGBColor(0xF8, 0xF6, 0xF0)
WARM        = RGBColor(0xE6, 0x9F, 0x00)
GREEN       = RGBColor(0x00, 0x9E, 0x73)
INK         = RGBColor(0x1F, 0x29, 0x37)
MUTED       = RGBColor(0x6B, 0x72, 0x80)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_NAVY  = RGBColor(0x60, 0x77, 0xC7)

HEAD_FONT = "Calibri"
BODY_FONT = "Calibri"


# ---------------------------------------------------------------------------
# Effort stats — fill with totals across local + SEAS cluster + Modal.
# Defaults below count LOCAL artifacts only; bump up to reflect full effort.
# ---------------------------------------------------------------------------
EFFORT = {
    "trajectories":      "500+",                          # local: 317
    "strategies":        "12",                            # exact
    "agent_steps":       "15K+",                          # local: 9,885
    "benchmarks":        "3",                             # SWE-bench Lite, τ-bench, long-doc
    "input_tokens":      "500M+",                         # local: 297M
    "phases":            "25+",                           # local: 20
    "api_cost":          "$200+",                         # local Haiku: ~$98
    "gpu_hours":         "100+",                          # rough total across Blackwell + SEAS + Modal
    "compute_summary":   "Anthropic Claude Haiku 4.5 (API)  ·  Qwen3-30B-A3B on Blackwell + SEAS cluster + Modal",
}


# ---------------------------------------------------------------------------
# Per-method result + run-count callouts.
# (run counts include local trajectories only; bump if you have seas/modal totals)
# ---------------------------------------------------------------------------
METHOD_RESULTS = {
    "none": {
        "result": "Pareto-dominates on N=10 SWE-bench Lite (Haiku 4.5)",
        "runs":   "84+ trajectories  ·  20 phases",
        "tone":   "neutral",
    },
    "naive_summary": {
        "result": "Always loses on cost — every fire writes uncached tokens",
        "runs":   "tested on τ-bench airline + SWE-bench Lite",
        "tone":   "negative",
    },
    "microcompact": {
        "result": "Rarely fires — needs an outlier-large observation",
        "runs":   "10 trajectories  ·  Phase B + C",
        "tone":   "neutral",
    },
    "prefix_preserving": {
        "result": "Stable cliff cost; loses ~30% on cost vs none",
        "runs":   "28 trajectories  ·  7 phases",
        "tone":   "negative",
    },
    "position_aware": {
        "result": "Proto-AdaptiveCache layout; no cost win at our scale",
        "runs":   "12 trajectories  ·  Phase A + C",
        "tone":   "neutral",
    },
    "evict_oldest": {
        "result": "FIFO baseline — drops irreplaceable old context too",
        "runs":   "6 trajectories  ·  Phase C",
        "tone":   "negative",
    },
    "smart_evict": {
        "result": "WINS on weak Qwen3 — prevents catastrophic OOC failures",
        "runs":   "59 trajectories  ·  14 phases",
        "tone":   "positive",
    },
    "score_periodic": {
        "result": "LLM scorer overhead exceeds bytes saved",
        "runs":   "limited tests  ·  Phase C v8",
        "tone":   "negative",
    },
    "llm_reorganizer": {
        "result": "Reorder + score — same overhead problem",
        "runs":   "30 trajectories  ·  9 phases",
        "tone":   "negative",
    },
    "consumption_evict_plain": {
        "result": "WINS pytest-7490 in 4 edits  ·  ties on retail chain (5/10)",
        "runs":   "Phase D 2× + Phase E (10 customers)",
        "tone":   "positive",
    },
    "consumption_evict_facts": {
        "result": "LOSES pytest-7490 — 18 edits, 0/2 success (mechanism finding)",
        "runs":   "Phase D 2×",
        "tone":   "negative",
    },
    "consumption_evict_outline": {
        "result": "Ties plain on retail chain (5/10) — cleanest design point",
        "runs":   "Phase E (10 customers)",
        "tone":   "positive",
    },
}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def add_blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def fill_background(slide, color):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height,
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    # Push to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, x, y, w, h, text, *,
             size=18, bold=False, italic=False, color=INK, font=BODY_FONT,
             align="left", valign="top"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM

    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}[align]
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_image_fit(slide, path, x, y, max_w, max_h):
    """Place image inside a max bounding box, preserving aspect ratio, centered."""
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = iw / ih
    box_ratio = max_w / max_h
    if ratio > box_ratio:
        # image wider than box → fit width
        w = max_w
        h = max_w / ratio
    else:
        h = max_h
        w = max_h * ratio
    cx = x + (max_w - w) / 2
    cy = y + (max_h - h) / 2
    return slide.shapes.add_picture(str(path), Inches(cx), Inches(cy),
                                     Inches(w), Inches(h))


def add_shape(slide, x, y, w, h, *, fill=None, line=None, line_w=0.5,
              shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def add_footer(slide, label, *, on_dark=False):
    color = LIGHT_NAVY if on_dark else MUTED
    add_text(slide, 0.5, SLIDE_H_IN - 0.4, 8, 0.3,
             "AdaptiveCache · Vlad Cainamisir · Harvard",
             size=10, color=color)
    add_text(slide, SLIDE_W_IN - 4, SLIDE_H_IN - 0.4, 3.5, 0.3,
             label, size=10, color=color, align="right")


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_title(prs):
    s = add_blank_slide(prs)
    fill_background(s, NAVY_DEEP)
    add_shape(s, 0, 0, 0.35, SLIDE_H_IN, fill=WARM)

    add_text(s, 1.2, 1.7, 11, 0.5,
             "ADAPTIVECACHE", size=14, bold=True, color=WARM, font=HEAD_FONT)
    add_text(s, 1.2, 2.25, 11, 2.4,
             "Compacting an AI agent's\nworking memory:\n12 ways to drop old data",
             size=42, bold=True, color=WHITE, font=HEAD_FONT)
    add_text(s, 1.2, 5.5, 11, 0.5,
             "What we learned from running 12 strategies on real coding and customer-service tasks.",
             size=17, italic=True, color=LIGHT_NAVY)
    add_text(s, 1.2, 6.7, 11, 0.4,
             "Vlad Cainamisir  ·  Harvard  ·  AdaptiveCache project",
             size=12, color=LIGHT_NAVY)


def slide_effort(prs, *, slide_num):
    """Header slide — what we actually ran."""
    s = add_blank_slide(prs)
    fill_background(s, NAVY_DEEP)
    add_shape(s, 0, 0, 0.35, SLIDE_H_IN, fill=WARM)

    add_text(s, 1.2, 0.7, 11, 0.5,
             "WHAT WE RAN", size=14, bold=True, color=WARM, font=HEAD_FONT)
    add_text(s, 1.2, 1.25, 11, 0.7,
             "Six months of compaction experiments, in numbers.",
             size=22, bold=True, color=WHITE, font=HEAD_FONT)

    # Big stat callouts in a 2 x 4 grid
    cells = [
        (EFFORT["trajectories"], "agent trajectories"),
        (EFFORT["strategies"],   "compaction strategies"),
        (EFFORT["agent_steps"],  "agent steps recorded"),
        (EFFORT["benchmarks"],   "benchmarks"),
        (EFFORT["input_tokens"], "input tokens"),
        (EFFORT["phases"],       "experimental phases"),
        (EFFORT["api_cost"],     "Anthropic API spend"),
        (EFFORT["gpu_hours"],    "GPU hours (Blackwell + SEAS + Modal)"),
    ]
    grid_x0, grid_y0 = 1.05, 2.55
    cw, ch = 2.78, 1.85
    gx, gy = 0.20, 0.2
    for i, (big, label) in enumerate(cells):
        col = i % 4
        row = i // 4
        x = grid_x0 + col * (cw + gx)
        y = grid_y0 + row * (ch + gy)
        add_shape(s, x, y, cw, ch, fill=NAVY)
        add_text(s, x, y + 0.1, cw, 0.95, big, size=36, bold=True,
                 color=WARM, align="center", valign="middle", font=HEAD_FONT)
        add_text(s, x + 0.1, y + 1.1, cw - 0.2, 0.65, label,
                 size=11, color=LIGHT_NAVY, align="center", valign="top")

    # Compute footer
    add_text(s, 1.2, 6.7, 11, 0.4,
             "Compute  ·  " + EFFORT["compute_summary"],
             size=11, italic=True, color=LIGHT_NAVY)

    add_footer(s, slide_num, on_dark=True)


def slide_problem(prs, *, slide_num):
    """Setup — what compaction is, no Post-it metaphor."""
    s = add_blank_slide(prs)
    fill_background(s, CREAM)

    add_text(s, 0.6, 0.4, 12, 0.6,
             "The problem: agent memory grows step by step.",
             size=26, bold=True, color=NAVY, font=HEAD_FONT)
    add_text(s, 0.6, 1.05, 12, 0.4,
             "Each tool call (read file, run test, search) adds 100s–1000s of tokens. Eventually we hit the model's limit.",
             size=14, italic=True, color=MUTED)

    # Three numbered problem statements as cards
    card_y, card_h, card_w, card_gap = 1.85, 2.55, 3.95, 0.25
    card_x0 = 0.6
    cards = [
        ("1", "Drop old observations",
         "Pick which past tool outputs to evict from the running context."),
        ("2", "Leave a placeholder",
         "Replace the dropped content with a short marker the agent can still see."),
        ("3", "Pay for the change",
         "Every rewrite invalidates the prefix cache → input bills uncached at 10× cost."),
    ]
    for i, (num, head, body) in enumerate(cards):
        x = card_x0 + i * (card_w + card_gap)
        add_shape(s, x, card_y, card_w, card_h, fill=WHITE,
                  line=NAVY, line_w=1.5)
        add_shape(s, x, card_y, card_w, 0.55, fill=NAVY)
        add_text(s, x + 0.25, card_y + 0.05, 0.6, 0.5, num,
                 size=20, bold=True, color=WARM, valign="middle")
        add_text(s, x + 0.95, card_y + 0.05, card_w - 1.1, 0.5, head,
                 size=14, bold=True, color=WHITE, valign="middle")
        add_text(s, x + 0.3, card_y + 0.85, card_w - 0.6, card_h - 1.0, body,
                 size=14, color=INK, valign="top")

    # Bottom callout: what we tested
    add_shape(s, 0.6, 4.85, 12.13, 0.7, fill=WARM)
    add_text(s, 0.8, 4.95, 12, 0.5,
             "We tested 12 ways of doing steps 1 + 2. The cost of step 3 decides whether they win.",
             size=15, bold=True, color=NAVY_DEEP, valign="middle")

    # Family-tree teaser image
    add_image_fit(s, FIG / "fig9_method_family_tree.png", 0.6, 5.7, 12.13, 1.5)

    add_footer(s, slide_num)


def slide_takeaway(prs, *, slide_num):
    """Closing slide — what we learned."""
    s = add_blank_slide(prs)
    fill_background(s, NAVY_DEEP)
    add_shape(s, 0, 0, 0.35, SLIDE_H_IN, fill=WARM)

    add_text(s, 1.2, 1.0, 11, 0.5,
             "WHAT WE LEARNED", size=14, bold=True, color=WARM, font=HEAD_FONT)
    add_text(s, 1.2, 1.7, 11, 1.4,
             "What you leave behind\nmatters more than what you remove.",
             size=34, bold=True, color=WHITE, font=HEAD_FONT)

    add_text(s, 1.2, 3.9, 11, 0.5,
             "1.  On strong agents, no compaction beats every heuristic — the cliff tax dominates.",
             size=17, color=LIGHT_NAVY)
    add_text(s, 1.2, 4.55, 11, 0.5,
             "2.  Drop-with-hole (no rewrite) is the only family that avoids the tax.",
             size=17, color=LIGHT_NAVY)
    add_text(s, 1.2, 5.2, 11, 0.5,
             "3.  Placeholder design dominates: minimal markers win, summaries anchor.",
             size=17, color=LIGHT_NAVY)
    add_text(s, 1.2, 5.85, 11, 0.5,
             "4.  Compaction wins on weak agents (Qwen3) — preventing context overflow.",
             size=17, color=LIGHT_NAVY)

    add_text(s, 1.2, 6.7, 11, 0.4,
             "Thank you.   ·   github.com/vladcainamisir/adaptivecache",
             size=12, color=LIGHT_NAVY)

    add_footer(s, slide_num, on_dark=True)


# ---------------------------------------------------------------------------
# Appendix slides
# ---------------------------------------------------------------------------

def appendix_divider(prs, title, subtitle):
    s = add_blank_slide(prs)
    fill_background(s, NAVY)
    add_text(s, 1.2, 2.8, 11, 1.5, title, size=44, bold=True, color=WHITE,
             font=HEAD_FONT)
    add_text(s, 1.2, 4.3, 11, 0.6, subtitle, size=20, italic=True,
             color=LIGHT_NAVY)


def slide_figure(prs, title, fig_path, caption, *, slide_num=None):
    """Generic content slide: title + dominant figure + 1-line caption."""
    s = add_blank_slide(prs)
    fill_background(s, CREAM)

    add_text(s, 0.6, 0.4, 12, 0.6, title, size=24, bold=True, color=NAVY,
             font=HEAD_FONT)
    if caption:
        add_text(s, 0.6, 1.0, 12, 0.4, caption, size=13, italic=True,
                 color=MUTED)

    add_image_fit(s, fig_path, 0.6, 1.55, 12.13, 5.4)

    if slide_num:
        add_footer(s, slide_num)


def slide_method(prs, fig_filename, family_label, family_color, *, slide_num):
    s = add_blank_slide(prs)
    fill_background(s, CREAM)

    # Family pill at top-left
    add_shape(s, 0.6, 0.45, 2.5, 0.45, fill=family_color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, 0.6, 0.45, 2.5, 0.45, family_label, size=11, bold=True,
             color=WHITE, align="center", valign="middle")

    # Method title
    method_name = fig_filename.replace(".png", "").replace("method_", "")
    method_name = method_name.split("_", 1)[1]  # strip "01_" prefix
    add_text(s, 3.3, 0.45, 9.5, 0.5, method_name, size=22, bold=True,
             color=NAVY, font=HEAD_FONT, valign="middle")

    # Result + runs callout band (key addition)
    info = METHOD_RESULTS.get(method_name, {})
    result_txt = info.get("result", "")
    runs_txt = info.get("runs", "")
    tone = info.get("tone", "neutral")
    band_fill = {"positive": GREEN, "negative": WARM, "neutral": NAVY}[tone]
    if result_txt:
        # Two-piece band: colored result strip + light "N runs" tag
        add_shape(s, 0.6, 1.05, 9.0, 0.6, fill=band_fill)
        add_text(s, 0.85, 1.05, 8.6, 0.6, result_txt,
                 size=14, bold=True, color=WHITE, valign="middle")
        # Runs box
        add_shape(s, 9.73, 1.05, 3.0, 0.6, fill=WHITE,
                  line=band_fill, line_w=1.25)
        add_text(s, 9.83, 1.05, 2.85, 0.6, runs_txt,
                 size=10, bold=True, color=band_fill, align="center", valign="middle")

    # Figure (slightly shorter to make room for the band)
    add_image_fit(s, FIG / fig_filename, 0.6, 1.85, 12.13, 5.0)

    add_footer(s, slide_num)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def _new_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def build_talk():
    """Method-walkthrough talk (~5 min)."""
    global prs
    prs = _new_pres()
    n_total = 17
    sn = lambda i: f"{i} / {n_total}"

    slide_title(prs)
    slide_effort(prs, slide_num=sn(2))
    slide_problem(prs, slide_num=sn(3))
    slide_figure(prs,
                 "The 12 strategies, grouped by family",
                 FIG / "fig9_method_family_tree.png",
                 "Six families. They differ on what signal they use to decide which observations to compact.",
                 slide_num=sn(4))

    method_specs = [
        ("method_01_none.png",                       "BASELINE",      MUTED),
        ("method_02_naive_summary.png",              "SUMMARIZE",     RGBColor(0x9D, 0x7C, 0xD8)),
        ("method_03_microcompact.png",               "SUMMARIZE",     RGBColor(0x9D, 0x7C, 0xD8)),
        ("method_04_prefix_preserving.png",          "SUMMARIZE",     RGBColor(0x9D, 0x7C, 0xD8)),
        ("method_05_position_aware.png",             "REORDER",       RGBColor(0x56, 0xB4, 0xE9)),
        ("method_06_evict_oldest.png",               "HEURISTIC",     WARM),
        ("method_07_smart_evict.png",                "HEURISTIC",     WARM),
        ("method_08_score_periodic.png",             "LLM-SCORED",    RGBColor(0xCC, 0x79, 0xA7)),
        ("method_09_llm_reorganizer.png",            "LLM-SCORED",    RGBColor(0xCC, 0x79, 0xA7)),
        ("method_10_consumption_evict_plain.png",    "ACTION-GRAPH",  RGBColor(0x00, 0x72, 0xB2)),
        ("method_11_consumption_evict_facts.png",    "ACTION-GRAPH",  RGBColor(0x00, 0x72, 0xB2)),
        ("method_12_consumption_evict_outline.png",  "ACTION-GRAPH",  RGBColor(0x00, 0x72, 0xB2)),
    ]
    for i, (fname, fam, col) in enumerate(method_specs, start=5):
        slide_method(prs, fname, fam, col, slide_num=sn(i))

    slide_takeaway(prs, slide_num=sn(n_total))

    prs.save(str(OUT_TALK))
    print(f"wrote {OUT_TALK}")
    print(f"  size: {OUT_TALK.stat().st_size / 1024:.0f} KB")
    print(f"  slides: {len(prs.slides)}")


def build_appendix():
    """Backup deck for Q&A — never opened during the talk itself."""
    global prs
    prs = _new_pres()
    n_total = 18  # divider + 3 framing + 12 methods + 3 supporting
    sn = lambda i: f"Backup {i} / {n_total}"

    appendix_divider(prs, "Appendix",
                     "Backup slides — open only if asked.")

    slide_figure(prs,
                 "The project — four phases, one negative-with-a-twist result",
                 FIG / "fig6_project_arc.png",
                 "Phase A: measurement → Phase B: τ-bench → Phase C: SWE-bench → Phase D: mechanism finding.",
                 slide_num=sn(2))

    slide_figure(prs,
                 "AdaptiveCache vision — context as a live, ordered memory",
                 FIG / "fig8_vision.png",
                 "Reorganize the prompt on compaction: pin stable items at the front; leave evictable holes in the suffix.",
                 slide_num=sn(3))

    slide_figure(prs,
                 "The compaction design space — 6 families, 12 methods",
                 FIG / "fig9_method_family_tree.png",
                 "All 12 policies we tested split on which signal they trust to decide what to drop.",
                 slide_num=sn(4))

    method_specs = [
        ("method_01_none.png",                       "BASELINE",      MUTED),
        ("method_02_naive_summary.png",              "SUMMARIZE",     RGBColor(0x9D, 0x7C, 0xD8)),
        ("method_03_microcompact.png",               "SUMMARIZE",     RGBColor(0x9D, 0x7C, 0xD8)),
        ("method_04_prefix_preserving.png",          "SUMMARIZE",     RGBColor(0x9D, 0x7C, 0xD8)),
        ("method_05_position_aware.png",             "REORDER",       RGBColor(0x56, 0xB4, 0xE9)),
        ("method_06_evict_oldest.png",               "HEURISTIC",     WARM),
        ("method_07_smart_evict.png",                "HEURISTIC",     WARM),
        ("method_08_score_periodic.png",             "LLM-SCORED",    RGBColor(0xCC, 0x79, 0xA7)),
        ("method_09_llm_reorganizer.png",            "LLM-SCORED",    RGBColor(0xCC, 0x79, 0xA7)),
        ("method_10_consumption_evict_plain.png",    "ACTION-GRAPH",  RGBColor(0x00, 0x72, 0xB2)),
        ("method_11_consumption_evict_facts.png",    "ACTION-GRAPH",  RGBColor(0x00, 0x72, 0xB2)),
        ("method_12_consumption_evict_outline.png",  "ACTION-GRAPH",  RGBColor(0x00, 0x72, 0xB2)),
    ]
    for i, (fname, fam, col) in enumerate(method_specs, start=5):
        slide_method(prs, fname, fam, col, slide_num=sn(i))

    slide_figure(prs,
                 "Headline result — none dominates on N=10 SWE-bench Lite",
                 FIG / "fig1_pareto_swebench.png",
                 "The training-free heuristic-compaction Pareto frontier vs. no-compaction baseline.",
                 slide_num=sn(17))

    slide_figure(prs,
                 "Cost decomposition — uncached input tokens dominate after cliffs",
                 FIG / "fig2_cost_decomposition.png",
                 "Each compaction event invalidates the prefix cache → uncached re-billing dominates.",
                 slide_num=sn(18))

    slide_figure(prs,
                 "Cliff tax — 10× cost amplification per compaction event",
                 FIG / "fig5_cliff_amplification.png",
                 "$0.10/MTok cached vs $1.00/MTok uncached. One cliff fires the difference for the rest of the trajectory.",
                 slide_num=sn(n_total))

    prs.save(str(OUT_APPENDIX))
    print(f"wrote {OUT_APPENDIX}")
    print(f"  size: {OUT_APPENDIX.stat().st_size / 1024:.0f} KB")
    print(f"  slides: {len(prs.slides)}")


def build():
    build_talk()
    build_appendix()


if __name__ == "__main__":
    build()

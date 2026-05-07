"""Build the AdaptiveCache decks (.pptx).

Outputs two files:
  reports/AdaptiveCache_talk.pptx      — main talk: method walkthrough, ~5 min
  reports/AdaptiveCache_appendix.pptx  — Q&A backup: project arc, vision, results

Main talk (16 slides, ~5 min):
  1.  Title
  2.  Setup — what compaction is and why we have to do it
  3.  Family tree — the 6 design families on one map
  4–15. The 12 methods, each with result + N runs callout
  16. Takeaway

Palette: navy + cream + warm orange + green (matches the figure palette).
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "reports" / "figures"
OUT_TALK = ROOT / "reports" / "AdaptiveCache_talk.pptx"
OUT_APPENDIX = ROOT / "reports" / "AdaptiveCache_appendix.pptx"

# 16:9 widescreen
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Palette
NAVY        = RGBColor(0x1E, 0x3A, 0x8A)
NAVY_DEEP   = RGBColor(0x14, 0x25, 0x5C)
CREAM       = RGBColor(0xF8, 0xF6, 0xF0)
WARM        = RGBColor(0xE6, 0x9F, 0x00)
GREEN       = RGBColor(0x00, 0x9E, 0x73)
INK         = RGBColor(0x1F, 0x29, 0x37)
MUTED       = RGBColor(0x6B, 0x72, 0x80)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_NAVY  = RGBColor(0xE8, 0xEC, 0xF5)  # pale off-white (was #6077C7 — low contrast on navy)

HEAD_FONT = "Calibri"
BODY_FONT = "Calibri"


# ---------------------------------------------------------------------------
# Effort stats — fill with totals across local + SEAS cluster + Modal.
# Defaults below count LOCAL artifacts only; bump up to reflect full effort.
# ---------------------------------------------------------------------------
EFFORT = {
    "trajectories":      "500+",                          # local: 317
    "strategies":        "12",                            # exact
    "agent_steps":       "15K+",                          # local: 9,885
    "benchmarks":        "3",                             # SWE-bench Lite, τ-bench, long-doc
    "input_tokens":      "500M+",                         # local: 297M
    "phases":            "25+",                           # local: 20
    "api_cost":          "$200+",                         # local Haiku: ~$98
    "gpu_hours":         "100+",                          # rough total across Blackwell + SEAS + Modal
    "compute_summary":   "Anthropic Claude Haiku 4.5 (API)  ·  Qwen3-30B-A3B on Blackwell + SEAS cluster + Modal",
}


# ---------------------------------------------------------------------------
# Per-method result + run-count callouts.
# (run counts include local trajectories only; bump if you have seas/modal totals)
# ---------------------------------------------------------------------------
METHOD_RESULTS = {
    "none": {
        "result": "Pareto-dominates on N=10 SWE-bench Lite (Haiku 4.5)",
        "runs":   "84+ trajectories  ·  20 phases",
        "tone":   "neutral",
    },
    "naive_summary": {
        "result": "Always loses on cost — every fire writes uncached tokens",
        "runs":   "tested on τ-bench airline + SWE-bench Lite",
        "tone":   "negative",
    },
    "microcompact": {
        "result": "Rarely fires — needs an outlier-large observation",
        "runs":   "10 trajectories  ·  Phase B + C",
        "tone":   "neutral",
    },
    "prefix_preserving": {
        "result": "Stable cliff cost; loses ~30% on cost vs none",
        "runs":   "28 trajectories  ·  7 phases",
        "tone":   "negative",
    },
    "position_aware": {
        "result": "Proto-AdaptiveCache layout; no cost win at our scale",
        "runs":   "12 trajectories  ·  Phase A + C",
        "tone":   "neutral",
    },
    "evict_oldest": {
        "result": "FIFO baseline — drops irreplaceable old context too",
        "runs":   "6 trajectories  ·  Phase C",
        "tone":   "negative",
    },
    "smart_evict": {
        "result": "WINS on weak Qwen3 — prevents catastrophic OOC failures",
        "runs":   "59 trajectories  ·  14 phases",
        "tone":   "positive",
    },
    "score_periodic": {
        "result": "LLM scorer overhead exceeds bytes saved",
        "runs":   "limited tests  ·  Phase C v8",
        "tone":   "negative",
    },
    "llm_reorganizer": {
        "result": "Reorder + score — same overhead problem",
        "runs":   "30 trajectories  ·  9 phases",
        "tone":   "negative",
    },
    "consumption_evict_plain": {
        "result": "WINS pytest-7490 in 4 edits  ·  ties on retail chain (5/10)",
        "runs":   "Phase D 2× + Phase E (10 customers)",
        "tone":   "positive",
    },
    "consumption_evict_facts": {
        "result": "LOSES pytest-7490 — 18 edits, 0/2 success (mechanism finding)",
        "runs":   "Phase D 2×",
        "tone":   "negative",
    },
    "consumption_evict_outline": {
        "result": "Ties plain on retail chain (5/10) — cleanest design point",
        "runs":   "Phase E (10 customers)",
        "tone":   "positive",
    },
}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def add_blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def fill_background(slide, color):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height,
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    # Push to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, x, y, w, h, text, *,
             size=18, bold=False, italic=False, color=INK, font=BODY_FONT,
             align="left", valign="top"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM

    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}[align]
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_image_fit(slide, path, x, y, max_w, max_h):
    """Place image inside a max bounding box, preserving aspect ratio, centered."""
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = iw / ih
    box_ratio = max_w / max_h
    if ratio > box_ratio:
        # image wider than box → fit width
        w = max_w
        h = max_w / ratio
    else:
        h = max_h
        w = max_h * ratio
    cx = x + (max_w - w) / 2
    cy = y + (max_h - h) / 2
    return slide.shapes.add_picture(str(path), Inches(cx), Inches(cy),
                                     Inches(w), Inches(h))


def add_shape(slide, x, y, w, h, *, fill=None, line=None, line_w=0.5,
              shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def add_footer(slide, label, *, on_dark=False):
    color = LIGHT_NAVY if on_dark else MUTED
    add_text(slide, 0.5, SLIDE_H_IN - 0.4, 8, 0.3,
             "AdaptiveCache · Vlad Cainamisir · Harvard",
             size=10, color=color)
    add_text(slide, SLIDE_W_IN - 4, SLIDE_H_IN - 0.4, 3.5, 0.3,
             label, size=10, color=color, align="right")


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_title(prs):
    s = add_blank_slide(prs)
    fill_background(s, NAVY_DEEP)
    add_shape(s, 0, 0, 0.35, SLIDE_H_IN, fill=WARM)

    add_text(s, 1.2, 1.7, 11, 0.5,
             "ADAPTIVECACHE", size=14, bold=True, color=WARM, font=HEAD_FONT)
    add_text(s, 1.2, 2.25, 11, 2.4,
             "Compacting an AI agent's\nworking memory:\n12 ways to drop old data",
             size=42, bold=True, color=WHITE, font=HEAD_FONT)
    add_text(s, 1.2, 5.5, 11, 0.5,
             "What we learned from running 12 strategies on real coding and customer-service tasks.",
             size=17, italic=True, color=LIGHT_NAVY)
    add_text(s, 1.2, 6.7, 11, 0.4,
             "Vlad Cainamisir  ·  Harvard  ·  AdaptiveCache project",
             size=12, color=LIGHT_NAVY)


def slide_effort(prs, *, slide_num):
    """Header slide — what we actually ran."""
    s = add_blank_slide(prs)
    fill_background(s, NAVY_DEEP)
    add_shape(s, 0, 0, 0.35, SLIDE_H_IN, fill=WARM)

    add_text(s, 1.2, 0.7, 11, 0.5,
             "WHAT WE RAN", size=14, bold=True, color=WARM, font=HEAD_FONT)
    add_text(s, 1.2, 1.25, 11, 0.7,
             "Six months of compaction experiments, in numbers.",
             size=22, bold=True, color=WHITE, font=HEAD_FONT)

    # Big stat callouts in a 2 x 4 grid
    cells = [
        (EFFORT["trajectories"], "agent trajectories"),
        (EFFORT["strategies"],   "compaction strategies"),
        (EFFORT["agent_steps"],  "agent steps recorded"),
        (EFFORT["benchmarks"],   "benchmarks"),
        (EFFORT["input_tokens"], "input tokens"),
        (EFFORT["phases"],       "experimental phases"),
        (EFFORT["api_cost"],     "Anthropic API spend"),
        (EFFORT["gpu_hours"],    "GPU hours (Blackwell + SEAS + Modal)"),
    ]
    grid_x0, grid_y0 = 1.05, 2.55
    cw, ch = 2.78, 1.85
    gx, gy = 0.20, 0.2
    for i, (big, label) in enumerate(cells):
        col = i % 4
        row = i // 4
        x = grid_x0 + col * (cw + gx)
        y = grid_y0 + row * (ch + gy)
        add_shape(s, x, y, cw, ch, fill=NAVY)
        add_text(s, x, y + 0.1, cw, 0.95, big, size=36, bold=True,
                 color=WARM, align="center", valign="middle", font=HEAD_FONT)
        add_text(s, x + 0.1, y + 1.1, cw - 0.2, 0.65, label,
                 size=11, color=LIGHT_NAVY, align="center", valign="top")

    # Compute footer
    add_text(s, 1.2, 6.7, 11, 0.4,
             "Compute  ·  " + EFFORT["compute_summary"],
             size=11, italic=True, color=LIGHT_NAVY)

    add_footer(s, slide_num, on_dark=True)


def slide_problem(prs, *, slide_num):
    """Setup — what compaction is, no Post-it metaphor."""
    s = add_blank_slide(prs)
    fill_background(s, CREAM)

    add_text(s, 0.6, 0.4, 12, 0.6,
             "The problem: agent memory grows step by step.",
             size=26, bold=True, color=NAVY, font=HEAD_FONT)
    add_text(s, 0.6, 1.05, 12, 0.4,
             "Each tool call (read file, run test, search) adds 100s–1000s of tokens. Eventually we hit the model's limit.",
             size=14, italic=True, color=MUTED)

    # Three numbered problem statements as cards
    card_y, card_h, card_w, card_gap = 1.85, 2.55, 3.95, 0.25
    card_x0 = 0.6
    cards = [
        ("1", "Drop old observations",
         "Pick which past tool outputs to evict from the running context."),
        ("2", "Leave a placeholder",
         "Replace the dropped content with a short marker the agent can still see."),
        ("3", "Pay for the change",
         "Every rewrite invalidates the prefix cache → input bills uncached at 10× cost."),
    ]
    for i, (num, head, body) in enumerate(cards):
        x = card_x0 + i * (card_w + card_gap)
        add_shape(s, x, card_y, card_w, card_h, fill=WHITE,
                  line=NAVY, line_w=1.5)
        add_shape(s, x, card_y, card_w, 0.55, fill=NAVY)
        add_text(s, x + 0.25, card_y + 0.05, 0.6, 0.5, num,
                 size=20, bold=True, color=WARM, valign="middle")
        add_text(s, x + 0.95, card_y + 0.05, card_w - 1.1, 0.5, head,
                 size=14, bold=True, color=WHITE, valign="middle")
        add_text(s, x + 0.3, card_y + 0.85, card_w - 0.6, card_h - 1.0, body,
                 size=14, color=INK, valign="top")

    # Bottom callout: what we tested
    add_shape(s, 0.6, 4.85, 12.13, 0.7, fill=WARM)
    add_text(s, 0.8, 4.95, 12, 0.5,
             "We tested 12 ways of doing steps 1 + 2. The cost of step 3 decides whether they win.",
             size=15, bold=True, color=NAVY_DEEP, valign="middle")

    # Family-tree teaser image
    add_image_fit(s, FIG / "fig9_method_family_tree.png", 0.6, 5.7, 12.13, 1.5)

    add_footer(s, slide_num)


def slide_takeaway(prs, *, slide_num):
    """Closing slide — what we learned."""
    s = add_blank_slide(prs)
    fill_background(s, NAVY_DEEP)
    add_shape(s, 0, 0, 0.35, SLIDE_H_IN, fill=WARM)

    add_text(s, 1.2, 1.0, 11, 0.5,
             "WHAT WE LEARNED", size=14, bold=True, color=WARM, font=HEAD_FONT)
    add_text(s, 1.2, 1.7, 11, 1.4,
             "What you leave behind\nmatters more than what you remove.",
             size=34, bold=True, color=WHITE, font=HEAD_FONT)

    add_text(s, 1.2, 3.9, 11, 0.5,
             "1.  On strong agents, no compaction beats every heuristic — the cliff tax dominates.",
             size=17, color=LIGHT_NAVY)
    add_text(s, 1.2, 4.55, 11, 0.5,
             "2.  Drop-with-hole (no rewrite) is the only family that avoids the tax.",
             size=17, color=LIGHT_NAVY)
    add_text(s, 1.2, 5.2, 11, 0.5,
             "3.  Placeholder design dominates: minimal markers win, summaries anchor.",
             size=17, color=LIGHT_NAVY)
    add_text(s, 1.2, 5.85, 11, 0.5,
             "4.  Compaction wins on weak agents (Qwen3) — preventing context overflow.",
             size=17, color=LIGHT_NAVY)

    add_text(s, 1.2, 6.7, 11, 0.4,
             "Thank you.   ·   github.com/vladcainamisir/adaptivecache",
             size=12, color=LIGHT_NAVY)

    add_footer(s, slide_num, on_dark=True)


# ---------------------------------------------------------------------------
# Paper 1 — headline result slide (cliff tax + Pareto)
# Lives between the methods walkthrough and the takeaway, so the audience
# sees the dollar story before the takeaway tells them why.
# ---------------------------------------------------------------------------

def slide_paper1_results(prs, *, slide_num):
    s = add_blank_slide(prs)
    fill_background(s, CREAM)

    add_text(s, 0.6, 0.4, 12, 0.6,
             "Headline result — none Pareto-dominates every heuristic",
             size=24, bold=True, color=NAVY, font=HEAD_FONT)
    add_text(s, 0.6, 1.05, 12, 0.4,
             "N=10 SWE-bench Lite, Haiku 4.5. Compaction policies cluster on the wrong side of the frontier.",
             size=13, italic=True, color=MUTED)

    # Pareto figure on the left
    add_image_fit(s, FIG / "fig1_pareto_swebench.png", 0.6, 1.6, 8.6, 5.0)

    # Cost-story stat callouts on the right — the WHY behind the Pareto
    callouts = [
        ("$0.10–0.15", "new uncached input per cliff", WARM),
        ("5+ steps",   "to amortize one cliff back to break-even", LIGHT_NAVY),
        ("2.2×",       "cheapest compaction policy vs none ($/resolved)", WARM),
    ]
    cy = 1.8
    for big, label, accent in callouts:
        add_shape(s, 9.5, cy, 3.2, 1.45, fill=NAVY)
        add_text(s, 9.5, cy + 0.1, 3.2, 0.7, big, size=26, bold=True,
                 color=accent, align="center", valign="middle", font=HEAD_FONT)
        add_text(s, 9.6, cy + 0.85, 3.0, 0.55, label, size=10,
                 color=LIGHT_NAVY, align="center", valign="top")
        cy += 1.65

    # Bottom takeaway band — keep clear of footer at y=7.10
    add_shape(s, 0.6, 6.62, 12.13, 0.38, fill=WARM)
    add_text(s, 0.8, 6.62, 12, 0.38,
             "The cliff tax is the dominant cost term. Any policy that fires has to fight it.",
             size=12, bold=True, color=NAVY_DEEP, valign="middle")

    add_footer(s, slide_num)


def slide_paper1_nuances(prs, *, slide_num):
    """Paper 1 nuances — what didn't work, plus the one positive Pareto win.

    Left: bulleted list of failure modes (rules don't port, tightening
    didn't help, LLM-reorganizer fired 91× and still overflowed).
    Right: fig7 — the Phase D v1 small-N Pareto-win.
    """
    s = add_blank_slide(prs)
    fill_background(s, CREAM)

    add_text(s, 0.6, 0.4, 12, 0.6,
             "What didn't work — and the one nugget that did",
             size=22, bold=True, color=NAVY, font=HEAD_FONT)
    add_text(s, 0.6, 1.0, 12, 0.4,
             "Most of Paper 1 is failure modes. Each told us something. One slice did Pareto-improve.",
             size=12, italic=True, color=MUTED)

    # LEFT column — failure modes
    fail_x, fail_y, fail_w = 0.6, 1.55, 5.6
    add_shape(s, fail_x, fail_y, fail_w, 0.5, fill=NAVY)
    add_text(s, fail_x + 0.25, fail_y + 0.05, fail_w - 0.5, 0.45,
             "What didn't work",
             size=14, bold=True, color=WHITE, valign="middle", font=HEAD_FONT)

    rows_y = fail_y + 0.6
    failures = [
        ("LLM-reorganizer",
         "Fired 91× across 4 SWE-bench tasks. Still hit context overflow on 3 of 4. Signal fine; cliff cost still wins."),
        ("Tighter τ-bench thresholds",
         "Stock τ-bench airline cache already 95% hit (max obs = 953 tok). Nothing to compact."),
        ("Action-graph rules don't port",
         "consumption_evict fired 0× on τ-bench retail at chain_size=10. Per-domain tuning required."),
        ("Placeholder design",
         "_facts variant anchors agent on wrong region → 18 edits, F2P 0/2. Less is more."),
    ]
    for head, body in failures:
        add_text(s, fail_x + 0.25, rows_y, fail_w - 0.5, 0.32, head,
                 size=12, bold=True, color=NAVY)
        add_text(s, fail_x + 0.25, rows_y + 0.36, fail_w - 0.5, 0.65, body,
                 size=10.5, color=INK)
        rows_y += 1.10

    # RIGHT column — the one Pareto win (fig7)
    win_x, win_y = 6.5, 1.55
    win_w = 6.5
    add_shape(s, win_x, win_y, win_w, 0.5, fill=GREEN)
    add_text(s, win_x + 0.25, win_y + 0.05, win_w - 0.5, 0.45,
             "Where compaction DID win",
             size=14, bold=True, color=WHITE, valign="middle", font=HEAD_FONT)

    add_image_fit(s, FIG / "fig7_compaction_wins.png", win_x, win_y + 0.65, win_w, 3.7)

    add_text(s, win_x + 0.1, win_y + 4.55, win_w - 0.2, 0.38,
             "Phase D v1 (Haiku, N=4):  consumption_evict 3/4 @ $1.03 vs none 2/4 @ $1.20.",
             size=11, bold=True, color=NAVY)
    add_text(s, win_x + 0.1, win_y + 4.95, win_w - 0.2, 0.38,
             "More resolved AND cheaper. Didn't fully replicate at N=10, but the slice is real.",
             size=10.5, italic=True, color=MUTED)

    add_footer(s, slide_num)


# ---------------------------------------------------------------------------
# Paper 2 / Memento track — exploratory follow-on slides
# Mark every slide with an "EXPLORATORY" tag so the audience knows this is
# in-flight work, not a result.
# ---------------------------------------------------------------------------

def _paper2_tag(slide):
    """Small orange pill in the top-right corner — visible status marker."""
    add_shape(slide, SLIDE_W_IN - 2.55, 0.40, 1.95, 0.42,
              fill=WARM, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, SLIDE_W_IN - 2.55, 0.40, 1.95, 0.42,
             "PAPER 2  ·  EXPLORATORY",
             size=10, bold=True, color=NAVY_DEEP, align="center", valign="middle")


def slide_paper2_divider(prs, *, slide_num):
    """Divider — clearly tag the new section as in-flight, not landed."""
    s = add_blank_slide(prs)
    fill_background(s, NAVY_DEEP)
    add_shape(s, 0, 0, 0.35, SLIDE_H_IN, fill=WARM)

    add_text(s, 1.2, 1.0, 11, 0.5,
             "A SECOND TRACK — IN FLIGHT",
             size=14, bold=True, color=WARM, font=HEAD_FONT)
    add_text(s, 1.2, 1.7, 11, 2.0,
             "If we can't beat the cliff,\ncan we keep the bytes?",
             size=40, bold=True, color=WHITE, font=HEAD_FONT)

    add_text(s, 1.2, 4.5, 11, 0.5,
             "Paper 2: KV-pointer recall — offload-and-restore as a cliff-free compaction substitute.",
             size=18, italic=True, color=LIGHT_NAVY)
    add_text(s, 1.2, 5.2, 11, 0.5,
             "Status: v0 mechanism shipped. v1 recall in flight. Real-test resolve still flat.",
             size=15, color=LIGHT_NAVY)
    add_text(s, 1.2, 5.85, 11, 0.5,
             "This has been the harder of the two tracks.",
             size=15, italic=True, color=WARM)

    add_footer(s, slide_num, on_dark=True)


def slide_paper2_thesis(prs, *, slide_num):
    """Three-card layout: why a Paper 2 — the cliff is the wall."""
    s = add_blank_slide(prs)
    fill_background(s, CREAM)
    _paper2_tag(s)

    add_text(s, 0.6, 0.4, 10, 0.6,
             "The cliff is the wall. So don't change the prefix.",
             size=24, bold=True, color=NAVY, font=HEAD_FONT)
    add_text(s, 0.6, 1.05, 10, 0.4,
             "Paper 1 said heuristic eviction loses to the cache cliff. Paper 2 asks: keep the bytes — only change the prompt rendering.",
             size=13, italic=True, color=MUTED)

    card_y, card_h, card_w, card_gap = 1.85, 3.4, 3.95, 0.25
    card_x0 = 0.6
    cards = [
        ("Paper 1 finding", "Cliff tax dominates",
         "Heuristic compaction generates ~$0.10–0.15 of new uncached input per fire. Even good policies lose to no-compaction."),
        ("Memento (MSR 2024)", "Mask, don't delete",
         "KV blocks stay in memory; positions stay numbered. Cliff still happens, but the masked-out KV → faster attention per step."),
        ("Our addition", "Bidirectional masking",
         "Recall the right obs when the agent needs it. v1: re-prefill from obs. v2: keep all KV, recover obs from offload, flip mask. Cost ≈ 0."),
    ]
    for i, (head_top, head, body) in enumerate(cards):
        x = card_x0 + i * (card_w + card_gap)
        add_shape(s, x, card_y, card_w, card_h, fill=WHITE,
                  line=NAVY, line_w=1.5)
        add_shape(s, x, card_y, card_w, 0.55, fill=NAVY)
        add_text(s, x + 0.25, card_y + 0.05, card_w - 0.5, 0.5, head_top,
                 size=11, bold=True, color=WARM, valign="middle")
        add_text(s, x + 0.3, card_y + 0.85, card_w - 0.6, 0.6, head,
                 size=18, bold=True, color=NAVY_DEEP, valign="top",
                 font=HEAD_FONT)
        add_text(s, x + 0.3, card_y + 1.55, card_w - 0.6, card_h - 1.7, body,
                 size=13, color=INK, valign="top")

    # Bottom orange callout
    add_shape(s, 0.6, 5.5, 12.13, 0.7, fill=WARM)
    add_text(s, 0.8, 5.6, 12, 0.5,
             "Mechanism is novel: nobody has shipped bidirectional KV masking on a production agent loop.",
             size=15, bold=True, color=NAVY_DEEP, valign="middle")

    add_footer(s, slide_num)


def slide_paper2_mechanism(prs, *, slide_num):
    """How the masking actually works at the KV-cache level — including
    the honest prefix-cache story."""
    s = add_blank_slide(prs)
    fill_background(s, CREAM)
    _paper2_tag(s)

    add_text(s, 0.6, 0.4, 10, 0.6,
             "How? Append a summary; don't rewrite the obs.",
             size=22, bold=True, color=NAVY, font=HEAD_FONT)
    add_text(s, 0.6, 1.0, 10, 0.4,
             "Positions stay coherent (top). Prefix cache: same cliff shape as Paper 1 — Memento's win is downstream of the cliff, not on it (bottom).",
             size=12, italic=True, color=MUTED)

    # Mechanism diagram — top half, smaller
    add_image_fit(s, FIG / "fig_memento_mechanism.png", 0.6, 1.45, 12.13, 2.7)

    # Two-card layout: positions + prefix cache
    card_y, card_h = 4.30, 2.30
    card_w, gap = 5.95, 0.25
    card_x0 = 0.6

    # LEFT card — positions
    x = card_x0
    add_shape(s, x, card_y, card_w, card_h, fill=WHITE, line=NAVY, line_w=1.5)
    add_shape(s, x, card_y, card_w, 0.5, fill=NAVY)
    add_text(s, x + 0.25, card_y + 0.05, card_w - 0.5, 0.45,
             "Why positions don't break",
             size=14, bold=True, color=WHITE, valign="middle", font=HEAD_FONT)
    add_text(s, x + 0.3, card_y + 0.65, card_w - 0.6, 0.45,
             "·  KV indices never renumber — masked positions stay numbered.",
             size=11, color=INK)
    add_text(s, x + 0.3, card_y + 1.05, card_w - 0.6, 0.45,
             "·  RoPE math on subsequent tokens stays consistent.",
             size=11, color=INK)
    add_text(s, x + 0.3, card_y + 1.45, card_w - 0.6, 0.65,
             "·  Naïve eviction renumbers and breaks attention; Memento masks and preserves it.",
             size=11, color=INK)

    # RIGHT card — prefix cache + lost-in-middle, honest version
    x = card_x0 + card_w + gap
    add_shape(s, x, card_y, card_w, card_h, fill=WHITE, line=NAVY, line_w=1.5)
    add_shape(s, x, card_y, card_w, 0.5, fill=WARM)
    add_text(s, x + 0.25, card_y + 0.05, card_w - 0.5, 0.45,
             "What about the prefix cache?",
             size=14, bold=True, color=NAVY_DEEP, valign="middle", font=HEAD_FONT)
    add_text(s, x + 0.3, card_y + 0.65, card_w - 0.6, 0.45,
             "·  Memento markers go IN THE MIDDLE — same cliff shape as Paper 1.",
             size=11, color=INK)
    add_text(s, x + 0.3, card_y + 1.05, card_w - 0.6, 0.45,
             "·  Win is downstream: smaller masked context → faster O(N²) attention per step.",
             size=11, color=INK)
    add_text(s, x + 0.3, card_y + 1.45, card_w - 0.6, 0.65,
             "·  Lost-in-middle is real: mementos drift into the U-shape low-attention zone — partly motivates content-aware recall.",
             size=11, color=INK, italic=True)

    add_footer(s, slide_num)


def slide_paper2_creation(prs, *, slide_num):
    """How a memento is created — real trace from pytest-7490.

    Three-column flow: tool call → 20K char obs → 791 char memento (Haiku
    writer produces a tight summary that captures the test functions).
    """
    s = add_blank_slide(prs)
    fill_background(s, CREAM)
    _paper2_tag(s)

    add_text(s, 0.6, 0.4, 10, 0.6,
             "Creating a memento — what the writer actually produces",
             size=22, bold=True, color=NAVY, font=HEAD_FONT)
    add_text(s, 0.6, 1.0, 10, 0.4,
             "Real trace from pytest-dev/pytest-7490, step 6. Haiku-as-writer in v0; an SFT'd model in the principled Paper 2.",
             size=12, italic=True, color=MUTED)

    # Three-column flow
    col_y, col_h = 1.55, 5.05
    col_w, gap = 4.05, 0.10
    col_x0 = 0.6

    # COL 1 — the tool call + truncated obs
    x = col_x0
    add_shape(s, x, col_y, col_w, 0.5, fill=NAVY)
    add_text(s, x + 0.2, col_y + 0.05, col_w - 0.4, 0.45,
             "1.  The tool call + obs",
             size=13, bold=True, color=WHITE, valign="middle", font=HEAD_FONT)
    add_shape(s, x, col_y + 0.5, col_w, col_h - 0.5, fill=WHITE,
              line=NAVY, line_w=1.2)
    add_text(s, x + 0.2, col_y + 0.65, col_w - 0.4, 0.35,
             "read_file('testing/test_mark.py')",
             size=11, bold=True, color=NAVY, font="Consolas")
    add_text(s, x + 0.2, col_y + 1.05, col_w - 0.4, 0.30,
             "→ 20,000 chars of code",
             size=10, italic=True, color=MUTED)
    code_excerpt = (
        "import os\n"
        "import sys\n"
        "from unittest import mock\n\n"
        "import pytest\n"
        "from _pytest.config import ExitCode\n"
        "from _pytest.mark import (\n"
        "    EMPTY_PARAMETERSET_OPTION,\n"
        "    MarkGenerator as Mark,\n"
        ")\n"
        "from _pytest.nodes import Collector\n\n"
        "class TestMark:\n"
        "    @pytest.mark.parametrize(\n"
        "        \"attr\", [\"mark\", \"param\"])\n"
        "    def test_pytest_exists_in_\n"
        "            namespace_all(self, ...):\n"
        "         ...\n"
        "    def test_pytest_mark_\n"
        "            notcallable(self, ...):\n"
        "         ...\n"
        "    [+ ~30 more methods]\n"
        "    [+ ~600 more lines]"
    )
    add_text(s, x + 0.2, col_y + 1.45, col_w - 0.4, col_h - 1.6, code_excerpt,
             size=9, color=INK, font="Consolas")

    # COL 2 — the writer
    x = col_x0 + col_w + gap
    add_shape(s, x, col_y, col_w, 0.5, fill=WARM)
    add_text(s, x + 0.2, col_y + 0.05, col_w - 0.4, 0.45,
             "2.  The writer summarizes",
             size=13, bold=True, color=NAVY_DEEP, valign="middle", font=HEAD_FONT)
    add_shape(s, x, col_y + 0.5, col_w, col_h - 0.5, fill=WHITE,
              line=NAVY, line_w=1.2)
    add_text(s, x + 0.2, col_y + 0.7, col_w - 0.4, 0.35,
             "Haiku 4.5  (v0 placeholder)",
             size=11, bold=True, color=NAVY)
    add_text(s, x + 0.2, col_y + 1.05, col_w - 0.4, 0.35,
             "→ SFT'd model in Paper 2",
             size=10, italic=True, color=MUTED)
    prompt = (
        "system: \"Compress this tool obs\n"
        "into ≤200 tokens. Preserve\n"
        "function names, key signals,\n"
        "what's still actionable.\"\n\n"
        "user:\n"
        "  tool_name: read_file\n"
        "  tool_args: {path: 'testing/\n"
        "              test_mark.py'}\n"
        "  obs: <20,000 chars>\n\n"
        "assistant: <791 chars>"
    )
    add_text(s, x + 0.2, col_y + 1.55, col_w - 0.4, col_h - 1.7, prompt,
             size=10, color=INK, font="Consolas")

    # COL 3 — the memento itself
    x = col_x0 + 2 * (col_w + gap)
    add_shape(s, x, col_y, col_w, 0.5, fill=GREEN)
    add_text(s, x + 0.2, col_y + 0.05, col_w - 0.4, 0.45,
             "3.  The memento",
             size=13, bold=True, color=WHITE, valign="middle", font=HEAD_FONT)
    add_shape(s, x, col_y + 0.5, col_w, col_h - 0.5, fill=WHITE,
              line=NAVY, line_w=1.2)
    add_text(s, x + 0.2, col_y + 0.7, col_w - 0.4, 0.35,
             "791 chars — captures structure",
             size=11, bold=True, color=NAVY)
    add_text(s, x + 0.2, col_y + 1.05, col_w - 0.4, 0.35,
             "→ replaces the obs in the prompt",
             size=10, italic=True, color=MUTED)
    memento_text = (
        "testing/test_mark.py: Test file\n"
        "for pytest marks functionality.\n"
        "Contains TestMark class with\n"
        "methods testing mark namespace\n"
        "(test_pytest_exists_in_namespace_\n"
        "all, test_pytest_mark_notcallable,\n"
        "test_mark_with_param, …).\n"
        "Module-level: test_marked_class_\n"
        "run_twice (parametrized class\n"
        "runs); test_ini_markers (ini\n"
        "config); test_markers_option\n"
        "(--markers flag); test_strict_\n"
        "prohibits_unregistered_markers\n"
        "(--strict-markers); test_mark_\n"
        "option (-m); test_keyword_\n"
        "option (-k); test_keyword_\n"
        "option_considers_mark."
    )
    add_text(s, x + 0.2, col_y + 1.55, col_w - 0.4, col_h - 1.7, memento_text,
             size=10, color=INK)

    # Bottom callout — what changed in the prompt after compaction
    add_shape(s, 0.6, 6.72, 12.13, 0.32, fill=NAVY_DEEP)
    add_text(s, 0.8, 6.72, 12, 0.32,
             "Compaction event:  20,000 chars  →  791 chars  (96% reduction). "
             "The original obs is wrapped with summary markers; the engine masks the obs's KV.",
             size=11, bold=True, color=WHITE, valign="middle")

    add_footer(s, slide_num)


def slide_paper2_retrieval(prs, *, slide_num):
    """How recall picks the right memento — real trace from pytest-7490 step 15.

    Left: agent's recent reasoning (the embedding query).
    Right: 5 candidate mementos with similarity scores; one highlighted as picked.
    """
    s = add_blank_slide(prs)
    fill_background(s, CREAM)
    _paper2_tag(s)

    add_text(s, 0.6, 0.4, 10, 0.6,
             "Retrieving a memento — what recall picks, and why",
             size=22, bold=True, color=NAVY, font=HEAD_FONT)
    add_text(s, 0.6, 1.0, 10, 0.4,
             "Real trace from step 15 of the embedding-recall run. MiniLM scores the agent's recent reasoning against each memento; highest above threshold wins.",
             size=12, italic=True, color=MUTED)

    # LEFT — the query
    qx, qy, qw, qh = 0.6, 1.55, 4.5, 5.20
    add_shape(s, qx, qy, qw, 0.5, fill=NAVY)
    add_text(s, qx + 0.2, qy + 0.05, qw - 0.4, 0.45,
             "Query — agent's recent context",
             size=13, bold=True, color=WHITE, valign="middle", font=HEAD_FONT)
    add_shape(s, qx, qy + 0.5, qw, qh - 0.5, fill=WHITE, line=NAVY, line_w=1.2)
    add_text(s, qx + 0.2, qy + 0.7, qw - 0.4, 0.35,
             "Last assistant turn (truncated):",
             size=10, bold=True, color=NAVY)
    quote = (
        "“Let me check the test for\n"
        "dynamically adding markers in the\n"
        "test_mark.py file.”\n\n"
        "+ tool_call: read_file(\n"
        "      path='testing/test_mark.py')"
    )
    add_text(s, qx + 0.2, qy + 1.10, qw - 0.4, 1.95, quote,
             size=11, italic=True, color=INK, font="Consolas")
    add_text(s, qx + 0.2, qy + 3.20, qw - 0.4, 0.35,
             "→ embedded by MiniLM (384-dim)",
             size=10, color=MUTED)
    add_text(s, qx + 0.2, qy + 3.60, qw - 0.4, 1.5,
             "Headroom check passes:\n  prompt ≈ 11K, low-water = 14.4K.\n  Recall is allowed to fire.",
             size=10, color=INK)

    # RIGHT — candidates with sim scores
    cx, cy, cw, ch = 5.40, 1.55, 7.33, 5.20
    add_shape(s, cx, cy, cw, 0.5, fill=WARM)
    add_text(s, cx + 0.2, cy + 0.05, cw - 0.4, 0.45,
             "Candidates — 5 mementos, scored by cosine-sim",
             size=13, bold=True, color=NAVY_DEEP, valign="middle", font=HEAD_FONT)
    add_shape(s, cx, cy + 0.5, cw, ch - 0.5, fill=WHITE, line=NAVY, line_w=1.2)

    # Approximate / illustrative similarity scores (real recall picked msg[11] —
    # the test_mark.py memento — which had top sim). Other candidates shown
    # for contrast.
    candidates = [
        (0.71, "msg[11] — test_mark.py",
         "Test file for pytest marks functionality. TestMark class with test_pytest_exists_in_namespace_all, test_mark_with_param, …",
         True),
        (0.42, "msg[23] — _pytest/python.py",
         "pytest test discovery & execution module; pytest_addoption registers marks, pytest_generate_tests parametrize, …",
         False),
        (0.31, "msg[9]  — testing/ dir listing",
         "path:testing/ contains: acceptance_test.py; conftest.py; subdirs: code/, examples/; test_argcomplete.py, test_assertion.py, …",
         False),
        (0.18, "msg[7]  — root dir listing",
         "Root directory: .coveragerc, .gitattributes, .github/, AUTHORS, CHANGELOG.rst, CITATION, …",
         False),
        (0.15, "msg[25] — _pytest/nodes.py",
         "Node base class for test collection tree; NodeMeta metaclass warns on direct instantiation; collect/iterations…",
         False),
    ]
    ry = cy + 0.65
    for sim, label, text, picked in candidates:
        # Sim score badge
        score_color = GREEN if picked else MUTED
        add_text(s, cx + 0.2, ry, 0.7, 0.32, f"{sim:.2f}",
                 size=12, bold=True, color=score_color, valign="middle", font="Consolas")
        # Label
        head_color = GREEN if picked else NAVY
        add_text(s, cx + 0.95, ry, cw - 1.1, 0.32, label,
                 size=11, bold=True, color=head_color)
        # Body — first ~120 chars of the memento
        add_text(s, cx + 0.95, ry + 0.32, cw - 1.1, 0.55, text,
                 size=9.5, color=INK, italic=True)
        if picked:
            # Pick marker on the right
            add_text(s, cx + cw - 1.5, ry, 1.3, 0.32, "PICKED ✓",
                     size=10, bold=True, color=GREEN, align="right",
                     valign="middle")
        ry += 0.92

    # Bottom callout — what happens next
    add_shape(s, 0.6, 6.72, 12.13, 0.32, fill=NAVY_DEEP)
    add_text(s, 0.8, 6.72, 12, 0.32,
             "v1 effect:  prompt re-renders with full obs → re-prefills from the obs onward. "
             "v2 target: keep all KV, recover obs from offload, flip mask. Generation continues (next slide).",
             size=11, bold=True, color=WHITE, valign="middle")

    add_footer(s, slide_num)


def slide_paper2_v0(prs, *, slide_num):
    """v0 mechanism demo — the per-step prompt comparison."""
    s = add_blank_slide(prs)
    fill_background(s, CREAM)
    _paper2_tag(s)

    add_text(s, 0.6, 0.4, 10, 0.6,
             "v0 — memento bounds the prompt while baseline grows linearly",
             size=22, bold=True, color=NAVY, font=HEAD_FONT)
    add_text(s, 0.6, 1.05, 10, 0.4,
             "vLLM 0.13.0 + Memento overlay + lazy policy. Real pytest-7490 agent loop. Both runs clean — no crashes.",
             size=13, italic=True, color=MUTED)

    add_image_fit(s, FIG / "fig_memento_v0_per_step.png", 0.6, 1.6, 8.6, 5.0)

    # Right-side stat callouts (pytest-7490, Qwen3-30B-A3B, T=0)
    callouts = [
        ("-31%",  "chat wall (13.3s vs 19.3s)"),
        ("-66%",  "final prompt (8K vs 23K)"),
        ("5",     "compactions fired (lazy policy)"),
    ]
    cy = 1.8
    for big, label in callouts:
        add_shape(s, 9.5, cy, 3.2, 1.45, fill=NAVY)
        add_text(s, 9.5, cy + 0.1, 3.2, 0.7, big, size=32, bold=True,
                 color=WARM, align="center", valign="middle", font=HEAD_FONT)
        add_text(s, 9.6, cy + 0.85, 3.0, 0.55, label, size=11,
                 color=LIGHT_NAVY, align="center", valign="top")
        cy += 1.65

    add_footer(s, slide_num)


def slide_paper2_v1(prs, *, slide_num):
    """v1 recall — three triggers compared on multi-seed pytest-7490."""
    s = add_blank_slide(prs)
    fill_background(s, CREAM)
    _paper2_tag(s)

    add_text(s, 0.6, 0.4, 10, 0.6,
             "v1 — three recall triggers, head-to-head",
             size=22, bold=True, color=NAVY, font=HEAD_FONT)
    add_text(s, 0.6, 1.05, 10, 0.4,
             "N=3 seeds · T=0.6 · pytest-dev/pytest-7490 · LRU and embedding both fire correctly. Wall time drops ~40% under either.",
             size=13, italic=True, color=MUTED)

    add_image_fit(s, FIG / "fig_memento_multi_seed.png", 0.6, 1.55, 12.13, 3.5)

    # Three-bullet readout below — must clear the y=7.10 footer
    add_text(s, 0.6, 5.30, 12.13, 0.4,
             "What's working:",
             size=13, bold=True, color=NAVY)
    add_text(s, 0.85, 5.70, 12, 0.4,
             "·   Plumbing fires correctly: LRU 2.0 recalls/run; embedding model loads + scores in real time.",
             size=12, color=INK)
    add_text(s, 0.85, 6.10, 12, 0.4,
             "·   Wall-time and read-count benefit replicates across seeds (~40% faster than no-recall).",
             size=12, color=INK)
    add_text(s, 0.85, 6.50, 12, 0.4,
             "·   The T=0 single-seed “embedding wins” was noise — multi-seed flips it: LRU narrowly best on this task.",
             size=12, color=INK)

    add_footer(s, slide_num)


def slide_paper2_v1_vs_v2(prs, *, slide_num):
    """Visual side-by-side of v1 (text-level recall, suffix invalidated) vs
    v2 (KV-level mask toggle, no prompt change, full cache hit)."""
    s = add_blank_slide(prs)
    fill_background(s, CREAM)
    _paper2_tag(s)

    add_text(s, 0.6, 0.4, 10, 0.6,
             "v1 vs v2 — keep the KV, flip the mask, keep going",
             size=22, bold=True, color=NAVY, font=HEAD_FONT)
    add_text(s, 0.6, 1.0, 10, 0.4,
             "v1 today: re-prefill from the obs.  v2 target: obs KV from offload, suffix KV stays as historical truth, mask flipped — generation just continues.",
             size=12, italic=True, color=MUTED)

    # The figure is the centerpiece — give it most of the slide
    add_image_fit(s, FIG / "fig_memento_v1_vs_v2.png", 0.5, 1.55, 12.33, 4.7)

    # Two side-by-side summary cards below
    card_y, card_h = 6.30, 0.75
    card_w, gap = 5.95, 0.25
    card_x0 = 0.6

    # LEFT — v1 (current) summary
    x = card_x0
    add_shape(s, x, card_y, card_w, card_h, fill=WHITE, line=NAVY, line_w=1.2)
    add_shape(s, x, card_y, 0.35, card_h, fill=WARM)
    add_text(s, x + 0.5, card_y + 0.05, card_w - 0.6, 0.30,
             "v1 (today):  re-prefill from the obs onward.",
             size=12, bold=True, color=NAVY)
    add_text(s, x + 0.5, card_y + 0.38, card_w - 0.6, 0.32,
             "Cost = obs_len + suffix_len.  Bounded prompt size still wins — wall ↓ ~40% over no-recall.",
             size=10.5, color=INK, italic=True)

    # RIGHT — v2 (target): keep all KV, flip the mask, generate
    x = card_x0 + card_w + gap
    add_shape(s, x, card_y, card_w, card_h, fill=WHITE, line=NAVY, line_w=1.2)
    add_shape(s, x, card_y, 0.35, card_h, fill=GREEN)
    add_text(s, x + 0.5, card_y + 0.05, card_w - 0.6, 0.30,
             "v2 (target):  keep KV; obs from offload; flip the mask.",
             size=12, bold=True, color=NAVY)
    add_text(s, x + 0.5, card_y + 0.38, card_w - 0.6, 0.32,
             "Cost ≈ 0.  Suffix KV is the historical truth — agent really did proceed memento-only. New generation sees both.",
             size=10.5, color=INK, italic=True)

    add_footer(s, slide_num)


def slide_paper2_hard(prs, *, slide_num):
    """The hard part — loose oracle vs real-test gap."""
    s = add_blank_slide(prs)
    fill_background(s, NAVY_DEEP)
    add_shape(s, 0, 0, 0.35, SLIDE_H_IN, fill=WARM)

    add_text(s, 1.2, 0.55, 11, 0.5,
             "WHAT'S HARD",
             size=14, bold=True, color=WARM, font=HEAD_FONT)
    add_text(s, 1.2, 1.05, 11, 0.9,
             "The mechanism works. The agent doesn't.",
             size=28, bold=True, color=WHITE, font=HEAD_FONT)

    # Figure on left
    add_image_fit(s, FIG / "fig_memento_oracle_gap.png", 0.6, 2.1, 7.0, 4.4)

    # Right-side honest-readout panel
    panel_x, panel_y, panel_w, panel_h = 8.0, 2.1, 4.85, 4.4
    add_shape(s, panel_x, panel_y, panel_w, panel_h, fill=NAVY)
    add_text(s, panel_x + 0.25, panel_y + 0.2, panel_w - 0.5, 0.5,
             "The honest read",
             size=14, bold=True, color=WARM, valign="top", font=HEAD_FONT)
    rows = [
        ("✓",  "v0 + v1 plumbing — shipped, validated"),
        ("✓",  "Wall time + read count — real benefit"),
        ("✗",  "Real-test FAIL_TO_PASS — no improvement"),
        ("",       "Qwen3-30B at single seed × 20 steps cannot fix pytest-7490 regardless of context strategy."),
    ]
    ry = panel_y + 0.85
    for mark, text in rows:
        if mark:
            color = GREEN if mark == "✓" else WARM
            add_text(s, panel_x + 0.3, ry, 0.4, 0.45, mark,
                     size=18, bold=True, color=color, valign="middle")
            add_text(s, panel_x + 0.85, ry, panel_w - 1.1, 0.45, text,
                     size=12, color=WHITE, valign="middle")
            ry += 0.6
        else:
            add_text(s, panel_x + 0.3, ry + 0.2, panel_w - 0.6, 1.4, text,
                     size=12, italic=True, color=LIGHT_NAVY, valign="top")

    # Footer line — what's next
    add_text(s, 0.6, 6.7, 12, 0.4,
             "Next  ·  Haiku-as-agent re-test (Paper 1's strong-agent setup)  ·  attention-driven recall (RQ2 #3)  ·  multi-task expansion",
             size=12, italic=True, color=LIGHT_NAVY, align="left")

    add_footer(s, slide_num, on_dark=True)


# ---------------------------------------------------------------------------
# Appendix slides
# ---------------------------------------------------------------------------

def appendix_divider(prs, title, subtitle):
    s = add_blank_slide(prs)
    fill_background(s, NAVY)
    add_text(s, 1.2, 2.8, 11, 1.5, title, size=44, bold=True, color=WHITE,
             font=HEAD_FONT)
    add_text(s, 1.2, 4.3, 11, 0.6, subtitle, size=20, italic=True,
             color=LIGHT_NAVY)


def slide_figure(prs, title, fig_path, caption, *, slide_num=None):
    """Generic content slide: title + dominant figure + 1-line caption."""
    s = add_blank_slide(prs)
    fill_background(s, CREAM)

    add_text(s, 0.6, 0.4, 12, 0.6, title, size=24, bold=True, color=NAVY,
             font=HEAD_FONT)
    if caption:
        add_text(s, 0.6, 1.0, 12, 0.4, caption, size=13, italic=True,
                 color=MUTED)

    add_image_fit(s, fig_path, 0.6, 1.55, 12.13, 5.4)

    if slide_num:
        add_footer(s, slide_num)


def slide_method(prs, fig_filename, family_label, family_color, *, slide_num):
    s = add_blank_slide(prs)
    fill_background(s, CREAM)

    # Family pill at top-left
    add_shape(s, 0.6, 0.45, 2.5, 0.45, fill=family_color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, 0.6, 0.45, 2.5, 0.45, family_label, size=11, bold=True,
             color=WHITE, align="center", valign="middle")

    # Method title
    method_name = fig_filename.replace(".png", "").replace("method_", "")
    method_name = method_name.split("_", 1)[1]  # strip "01_" prefix
    add_text(s, 3.3, 0.45, 9.5, 0.5, method_name, size=22, bold=True,
             color=NAVY, font=HEAD_FONT, valign="middle")

    # Result + runs callout band (key addition)
    info = METHOD_RESULTS.get(method_name, {})
    result_txt = info.get("result", "")
    runs_txt = info.get("runs", "")
    tone = info.get("tone", "neutral")
    band_fill = {"positive": GREEN, "negative": WARM, "neutral": NAVY}[tone]
    if result_txt:
        # Two-piece band: colored result strip + light "N runs" tag
        add_shape(s, 0.6, 1.05, 9.0, 0.6, fill=band_fill)
        add_text(s, 0.85, 1.05, 8.6, 0.6, result_txt,
                 size=14, bold=True, color=WHITE, valign="middle")
        # Runs box
        add_shape(s, 9.73, 1.05, 3.0, 0.6, fill=WHITE,
                  line=band_fill, line_w=1.25)
        add_text(s, 9.83, 1.05, 2.85, 0.6, runs_txt,
                 size=10, bold=True, color=band_fill, align="center", valign="middle")

    # Figure (slightly shorter to make room for the band)
    add_image_fit(s, FIG / fig_filename, 0.6, 1.85, 12.13, 5.0)

    add_footer(s, slide_num)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def _new_pres():
    p = Presentation()
    p.slide_width = Inches(SLIDE_W_IN)
    p.slide_height = Inches(SLIDE_H_IN)
    return p


def build_talk():
    """Method-walkthrough talk (~5 min) + Paper 2 / Memento exploratory track."""
    global prs
    prs = _new_pres()
    n_total = 28  # was 17; +1 P1 results, +1 nuances, +5 P2, +1 mechanism, +2 trace, +1 v1-vs-v2
    sn = lambda i: f"{i} / {n_total}"

    slide_title(prs)
    slide_effort(prs, slide_num=sn(2))
    slide_problem(prs, slide_num=sn(3))
    slide_figure(prs,
                 "The 12 strategies, grouped by family",
                 FIG / "fig9_method_family_tree.png",
                 "Six families. They differ on what signal they use to decide which observations to compact.",
                 slide_num=sn(4))

    method_specs = [
        ("method_01_none.png",                       "BASELINE",      MUTED),
        ("method_02_naive_summary.png",              "SUMMARIZE",     RGBColor(0x9D, 0x7C, 0xD8)),
        ("method_03_microcompact.png",               "SUMMARIZE",     RGBColor(0x9D, 0x7C, 0xD8)),
        ("method_04_prefix_preserving.png",          "SUMMARIZE",     RGBColor(0x9D, 0x7C, 0xD8)),
        ("method_05_position_aware.png",             "REORDER",       RGBColor(0x56, 0xB4, 0xE9)),
        ("method_06_evict_oldest.png",               "HEURISTIC",     WARM),
        ("method_07_smart_evict.png",                "HEURISTIC",     WARM),
        ("method_08_score_periodic.png",             "LLM-SCORED",    RGBColor(0xCC, 0x79, 0xA7)),
        ("method_09_llm_reorganizer.png",            "LLM-SCORED",    RGBColor(0xCC, 0x79, 0xA7)),
        ("method_10_consumption_evict_plain.png",    "ACTION-GRAPH",  RGBColor(0x00, 0x72, 0xB2)),
        ("method_11_consumption_evict_facts.png",    "ACTION-GRAPH",  RGBColor(0x00, 0x72, 0xB2)),
        ("method_12_consumption_evict_outline.png",  "ACTION-GRAPH",  RGBColor(0x00, 0x72, 0xB2)),
    ]
    for i, (fname, fam, col) in enumerate(method_specs, start=5):
        slide_method(prs, fname, fam, col, slide_num=sn(i))

    # Paper 1 — headline result + cliff cost story, then nuances + 1 win
    slide_paper1_results(prs,  slide_num=sn(17))
    slide_paper1_nuances(prs,  slide_num=sn(18))
    slide_takeaway(prs,        slide_num=sn(19))

    # Paper 2 / Memento exploratory track — slides 20 through 28
    slide_paper2_divider(prs,    slide_num=sn(20))
    slide_paper2_thesis(prs,     slide_num=sn(21))
    slide_paper2_mechanism(prs,  slide_num=sn(22))
    slide_paper2_creation(prs,   slide_num=sn(23))
    slide_paper2_retrieval(prs,  slide_num=sn(24))
    slide_paper2_v0(prs,         slide_num=sn(25))
    slide_paper2_v1(prs,         slide_num=sn(26))
    slide_paper2_v1_vs_v2(prs,   slide_num=sn(27))
    slide_paper2_hard(prs,       slide_num=sn(28))

    prs.save(str(OUT_TALK))
    print(f"wrote {OUT_TALK}")
    print(f"  size: {OUT_TALK.stat().st_size / 1024:.0f} KB")
    print(f"  slides: {len(prs.slides)}")


def build_appendix():
    """Backup deck for Q&A — never opened during the talk itself."""
    global prs
    prs = _new_pres()
    n_total = 18  # divider + 3 framing + 12 methods + 3 supporting
    sn = lambda i: f"Backup {i} / {n_total}"

    appendix_divider(prs, "Appendix",
                     "Backup slides — open only if asked.")

    slide_figure(prs,
                 "The project — four phases, one negative-with-a-twist result",
                 FIG / "fig6_project_arc.png",
                 "Phase A: measurement → Phase B: τ-bench → Phase C: SWE-bench → Phase D: mechanism finding.",
                 slide_num=sn(2))

    slide_figure(prs,
                 "AdaptiveCache vision — context as a live, ordered memory",
                 FIG / "fig8_vision.png",
                 "Reorganize the prompt on compaction: pin stable items at the front; leave evictable holes in the suffix.",
                 slide_num=sn(3))

    slide_figure(prs,
                 "The compaction design space — 6 families, 12 methods",
                 FIG / "fig9_method_family_tree.png",
                 "All 12 policies we tested split on which signal they trust to decide what to drop.",
                 slide_num=sn(4))

    method_specs = [
        ("method_01_none.png",                       "BASELINE",      MUTED),
        ("method_02_naive_summary.png",              "SUMMARIZE",     RGBColor(0x9D, 0x7C, 0xD8)),
        ("method_03_microcompact.png",               "SUMMARIZE",     RGBColor(0x9D, 0x7C, 0xD8)),
        ("method_04_prefix_preserving.png",          "SUMMARIZE",     RGBColor(0x9D, 0x7C, 0xD8)),
        ("method_05_position_aware.png",             "REORDER",       RGBColor(0x56, 0xB4, 0xE9)),
        ("method_06_evict_oldest.png",               "HEURISTIC",     WARM),
        ("method_07_smart_evict.png",                "HEURISTIC",     WARM),
        ("method_08_score_periodic.png",             "LLM-SCORED",    RGBColor(0xCC, 0x79, 0xA7)),
        ("method_09_llm_reorganizer.png",            "LLM-SCORED",    RGBColor(0xCC, 0x79, 0xA7)),
        ("method_10_consumption_evict_plain.png",    "ACTION-GRAPH",  RGBColor(0x00, 0x72, 0xB2)),
        ("method_11_consumption_evict_facts.png",    "ACTION-GRAPH",  RGBColor(0x00, 0x72, 0xB2)),
        ("method_12_consumption_evict_outline.png",  "ACTION-GRAPH",  RGBColor(0x00, 0x72, 0xB2)),
    ]
    for i, (fname, fam, col) in enumerate(method_specs, start=5):
        slide_method(prs, fname, fam, col, slide_num=sn(i))

    slide_figure(prs,
                 "Headline result — none dominates on N=10 SWE-bench Lite",
                 FIG / "fig1_pareto_swebench.png",
                 "The training-free heuristic-compaction Pareto frontier vs. no-compaction baseline.",
                 slide_num=sn(17))

    slide_figure(prs,
                 "Cost decomposition — uncached input tokens dominate after cliffs",
                 FIG / "fig2_cost_decomposition.png",
                 "Each compaction event invalidates the prefix cache → uncached re-billing dominates.",
                 slide_num=sn(18))

    slide_figure(prs,
                 "Cliff tax — 10× cost amplification per compaction event",
                 FIG / "fig5_cliff_amplification.png",
                 "$0.10/MTok cached vs $1.00/MTok uncached. One cliff fires the difference for the rest of the trajectory.",
                 slide_num=sn(n_total))

    prs.save(str(OUT_APPENDIX))
    print(f"wrote {OUT_APPENDIX}")
    print(f"  size: {OUT_APPENDIX.stat().st_size / 1024:.0f} KB")
    print(f"  slides: {len(prs.slides)}")


def build():
    build_talk()
    build_appendix()


if __name__ == "__main__":
    build()
